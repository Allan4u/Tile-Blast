# -*- coding: utf-8 -*-
"""
Generator presentasi PPTX untuk project Android TileBlast.

Tema: dark background (#0A0A0A), accent gold (#FFD700), text putih.
Aspek 16:9. ~26 slide.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ---------------------------------------------------------------------------
# Konstanta tema
# ---------------------------------------------------------------------------
BG_COLOR = RGBColor(0x0A, 0x0A, 0x0A)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
WHITE = RGBColor(0xF5, 0xF5, 0xF5)
GREY = RGBColor(0xAA, 0xAA, 0xAA)
ACCENT_BLUE = RGBColor(0x42, 0xA5, 0xF5)

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)

OUTPUT = Path(__file__).parent / "TileBlast_Presentation.pptx"


def set_bg(slide, color=BG_COLOR):
    """Isi background slide dengan warna solid."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height,
                 text, *, size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 font_name="Calibri"):
    """Tambah text box berisi satu baris teks."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def add_bullets(slide, left, top, width, height, bullets,
                *, size=18, color=WHITE, bullet_color=GOLD,
                line_spacing=1.15, bullet_char="•"):
    """Tambah text box berisi list bullet."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    for i, item in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # Bullet character run (gold)
        bullet_run = p.add_run()
        bullet_run.text = f"{bullet_char}  "
        bullet_run.font.size = Pt(size)
        bullet_run.font.bold = True
        bullet_run.font.color.rgb = bullet_color
        bullet_run.font.name = "Calibri"
        # Content run (white)
        content_run = p.add_run()
        content_run.text = item
        content_run.font.size = Pt(size)
        content_run.font.color.rgb = color
        content_run.font.name = "Calibri"
    return tb


def add_title(slide, title_text, *, size=32, color=GOLD):
    """Tambah judul slide standar."""
    return add_text_box(
        slide, Inches(0.5), Inches(0.35),
        Inches(12.3), Inches(0.9),
        title_text, size=size, bold=True, color=color
    )


def add_subtitle_bar(slide, subtitle):
    """Tambah garis horizontal gold tipis di bawah judul."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.25),
        Inches(2.0), Inches(0.05),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()
    if subtitle:
        add_text_box(
            slide, Inches(0.5), Inches(1.35),
            Inches(12.3), Inches(0.4),
            subtitle, size=14, color=GREY, bold=False,
        )


def add_footer(slide, page_num, total_pages):
    """Tambah footer kecil di pojok kanan bawah."""
    add_text_box(
        slide, Inches(11.5), Inches(7.05),
        Inches(1.7), Inches(0.35),
        f"TileBlast  •  {page_num}/{total_pages}",
        size=10, color=GREY, align=PP_ALIGN.RIGHT,
    )


def make_content_slide(prs, title, subtitle, bullets, *,
                       bullet_size=18, page_num=1, total=1):
    """Helper: bikin slide dark + judul gold + bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_bg(slide)
    add_title(slide, title)
    add_subtitle_bar(slide, subtitle)
    add_bullets(
        slide, Inches(0.5), Inches(1.85),
        Inches(12.3), Inches(5.0),
        bullets, size=bullet_size,
    )
    add_footer(slide, page_num, total)
    return slide


def make_two_col_slide(prs, title, subtitle, left_title, left_bullets,
                       right_title, right_bullets, *, page_num=1, total=1):
    """Helper: dua kolom, masing-masing punya sub-judul + bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, title)
    add_subtitle_bar(slide, subtitle)

    # Kolom kiri
    add_text_box(
        slide, Inches(0.5), Inches(1.85),
        Inches(6.0), Inches(0.5),
        left_title, size=18, bold=True, color=GOLD,
    )
    add_bullets(
        slide, Inches(0.5), Inches(2.4),
        Inches(6.0), Inches(4.5),
        left_bullets, size=15,
    )

    # Kolom kanan
    add_text_box(
        slide, Inches(6.85), Inches(1.85),
        Inches(6.0), Inches(0.5),
        right_title, size=18, bold=True, color=GOLD,
    )
    add_bullets(
        slide, Inches(6.85), Inches(2.4),
        Inches(6.0), Inches(4.5),
        right_bullets, size=15,
    )

    add_footer(slide, page_num, total)
    return slide


# ---------------------------------------------------------------------------
# Slide builder
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Total slide direncanakan = 26
    TOTAL = 26

    # ------------------------------------------------------------------
    # Slide 1 — Cover
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    # Garis gold besar di kiri
    bar = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(2.5),
        Inches(0.18), Inches(2.5),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    add_text_box(
        s, Inches(0.85), Inches(2.4),
        Inches(11.5), Inches(1.4),
        "TileBlast",
        size=80, bold=True, color=GOLD, font_name="Calibri",
    )
    add_text_box(
        s, Inches(0.85), Inches(3.7),
        Inches(11.5), Inches(0.7),
        "Block Puzzle Game with 10 Major Features",
        size=26, color=WHITE, font_name="Calibri",
    )
    add_text_box(
        s, Inches(0.85), Inches(4.5),
        Inches(11.5), Inches(0.6),
        "Android  •  Java 17  •  Custom Canvas Rendering  •  Firebase",
        size=16, color=GREY,
    )
    add_text_box(
        s, Inches(0.85), Inches(6.4),
        Inches(11.5), Inches(0.4),
        "Presented by: [Nama Developer]",
        size=14, color=WHITE,
    )
    add_text_box(
        s, Inches(0.85), Inches(6.8),
        Inches(11.5), Inches(0.4),
        "Tanggal: [Tanggal Presentasi]",
        size=14, color=GREY,
    )

    # ------------------------------------------------------------------
    # Slide 2 — Daftar Isi
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Daftar Isi",
        subtitle="Apa yang akan kita bahas hari ini",
        bullets=[
            "Tentang TileBlast — overview project",
            "Tech Stack & Build Configuration",
            "Arsitektur 3-Layer & Modul",
            "Game Modes — 5 mode permainan",
            "10 Fitur Major (slide tersendiri per fitur)",
            "Visual Rendering, Audio System, Persistence",
            "Property-Based Testing dengan jqwik",
            "Spec-Driven Development & Statistik Project",
            "Tantangan Teknis & Future Work",
            "Demo Highlights, Penutup & Q&A",
        ],
        bullet_size=20,
        page_num=2, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 3 — Tentang TileBlast
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Tentang TileBlast",
        subtitle="Block puzzle modern untuk Android",
        bullets=[
            "Game block puzzle untuk Android, terinspirasi dari Blockudoku dan 1010!",
            "Tarik & lepas piece ke board, isi baris atau kolom penuh untuk dibersihkan",
            "Custom rendering 100% via Android Canvas — tanpa engine pihak ketiga",
            "Fokus pada juiciness: particle effects, screen flash, zoom pulse, combo text",
            "10 fitur utama: dari power-ups, special tiles, daily challenge, sampai online leaderboard",
            "Package root: com.allan.tileblast — ditulis dalam Java 17 modular",
        ],
        page_num=3, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 4 — Tech Stack
    # ------------------------------------------------------------------
    make_two_col_slide(
        prs,
        title="Tech Stack",
        subtitle="Komponen teknologi yang dipakai",
        left_title="Platform & Build",
        left_bullets=[
            "Android SDK 34 (compileSdk), minSdk 24 — Android 7.0+",
            "Java 17 (sourceCompatibility & targetCompatibility)",
            "Gradle Kotlin DSL + AGP 8.4.0",
            "Core library desugaring untuk java.time di minSdk 24",
            "Material Components 1.11.0 + AndroidX AppCompat",
        ],
        right_title="Backend & Testing",
        right_bullets=[
            "Firebase BoM 34.13.0 (Auth + Firestore)",
            "Google Play Services Auth 21.3.0 (Google Sign-In)",
            "SoundPool untuk efek audio low-latency",
            "Custom Canvas View — tanpa OpenGL/Compose",
            "JUnit 5 + jqwik 1.8.5 untuk property-based testing",
            "Mockito 5.14.2 untuk unit test biasa",
        ],
        page_num=4, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 5 — Arsitektur
    # ------------------------------------------------------------------
    make_two_col_slide(
        prs,
        title="Arsitektur 3-Layer",
        subtitle="Pemisahan UI, Game Logic, dan Persistence",
        left_title="Layer Hierarchy",
        left_bullets=[
            "UI Layer: Activities + custom View (GameView)",
            "Logic Layer: Manager classes pure JVM",
            "Persistence Layer: SharedPreferences + Bundle + Firestore",
            "Singletons: AppContext, AuthManager, LeaderboardService, ThemeManager, AudioManager",
            "Komunikasi via callback interface (GameView.GameCallback)",
        ],
        right_title="Modul (package)",
        right_bullets=[
            "game/  — Board, Hand, Piece, GameView, ScoreManager, …",
            "game/effects/  — Particle, ScreenFlash, ZoomPulse, …",
            "audio/  — AudioManager (SoundPool wrapper)",
            "storage/  — StorageManager (SharedPreferences)",
            "leaderboard/  — Auth, Service, Adapter, ScoreQueue",
            "daily/  — DailyChallengeEngine, ChallengeStorage, StreakTracker, CalendarView",
            "stats/, progression/, theme/  — fitur lain",
        ],
        page_num=5, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 6 — Game Modes
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Game Modes",
        subtitle="5 mode permainan, 1 game core",
        bullets=[
            "Classic — board 8×8, hand 3 piece, no timer (mode default)",
            "Chaos — board 10×10, hand 5 piece, special tiles aktif",
            "Timed 60 — board 8×8, durasi 60 detik, scoring multiplier",
            "Timed 90 — board 8×8, durasi 90 detik, multiplier sampai 3.0×",
            "Daily Challenge — seed YYYYMMDD, target 2000, kalender 30 hari + streak",
            "Semua mode share GameView yang sama, dibedakan via parameter setup()",
        ],
        page_num=6, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 7 — Fitur 1: Enhanced Visual Feedback
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Fitur 1 — Enhanced Visual Feedback",
        subtitle="Juiciness via particles, flash, zoom, combo text",
        bullets=[
            "EffectManager — orchestrator semua efek visual; clamp deltaMs ≤ 33ms",
            "ParticleSystem + ParticlePool — pool 512 partikel reusable, gravity 800 px/s²",
            "Line break: 4–8 partikel/cell, lifetime 400–700ms; Perfect Clear: 80–120 partikel emas",
            "ScreenFlash — combo ≥3 putih 15%, combo 4+ emas opasitas 0.15+0.05·(level-3) cap 0.40",
            "ZoomPulse — combo ≥4 atau lines ≥3, scale 1.0 → 1.03 → 1.0 dalam 300ms",
            "ComboTextAnim — overshoot scale-in → hold → fade; PerfectClear pakai durasi 800ms",
            "ScorePopManager — pool 8 slot, stacking ≤200ms window, rise 60dp + fade",
        ],
        page_num=7, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 8 — Fitur 2: Ghost Preview & Hints
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Fitur 2 — Ghost Preview & Hints",
        subtitle="Bantu pemain melihat placement optimal",
        bullets=[
            "Ghost preview — Board.setHover() menandai cell HOVERED + HOVERED_BREAK_*",
            "Pulse alpha pada ghost: berdenyut antara min–max via computePulseAlpha()",
            "Highlight emas line yang akan break (HOVERED_BREAK_FILLED / HOVERED_BREAK_EMPTY)",
            "Hint button — 3 hint per game, decrement on tap",
            "HintCalculator scoring: lines·1000 + adjacent·10 + edge·5; tiebreak: distance ke center",
            "Pure JVM class — bisa di-test tanpa Android framework dependencies",
        ],
        page_num=8, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 9 — Fitur 3: Power-Ups System
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Fitur 3 — Power-Ups System",
        subtitle="4 tipe, state machine, snapshot deep-copy",
        bullets=[
            "4 tipe: BOMB (3×3 clear, +10 per filled), LINE_SWEEP (row/column), ROTATE (90° CW), UNDO",
            "Activation state machine: IDLE → TARGETING → EXECUTING → IDLE",
            "Inventory cap 2 per type (MAX_PER_TYPE); surplus dari grant di-discard",
            "Acquisition: grantOnCombo(level≥4) + grantOnScore (milestones 1k/5k/10k/25k/50k)",
            "BOMB & LINE_SWEEP boleh saat gameOver; ROTATE wajib saat dragging; UNDO butuh snapshot bersih",
            "BoardSnapshot.capture() — deep-copy cells, hand matrices (preserve rotation), score, combo",
            "GameContext immutable struct untuk gating tanpa back-ref ke GameView",
        ],
        page_num=9, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 10 — Fitur 4: Special Tiles
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Fitur 4 — Special Tiles",
        subtitle="Frozen, Locked, Bomb — hanya di Chaos & Daily",
        bullets=[
            "FROZEN (cap 5, 10%) → 2× clear: FROZEN → FROZEN_CRACKED → EMPTY",
            "LOCKED (cap 3, 5%) → 1× clear, tidak bisa ditimpa piece",
            "BOMB_TILE (cap 2, 3%) → countdown 5 turn → 3×3 detonation bounded by edges",
            "Chain bomb: jika detonation menyentuh BOMB lain → onChainBombGameOver()",
            "Spawn turn-gated: FROZEN minimal turn ≥ 6 (FROZEN_MIN_TURN=5, strict greater)",
            "Tiap turn: tiga roll independen (FROZEN, LOCKED, BOMB) — bisa spawn ≤3 sekaligus",
            "Class: SpecialTilesManager owns evaluateSpawns(), decrementBombsAndDetonate()",
        ],
        page_num=10, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 11 — Fitur 5: Daily Challenge
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Fitur 5 — Daily Challenge",
        subtitle="Puzzle harian deterministik + streak system",
        bullets=[
            "Seed deterministik: year×10000 + month×100 + day (YYYYMMDD long)",
            "Random rng = new Random(seed) → semua player dapet puzzle identik per tanggal",
            "Board 8×8, hand 3, target 2000 untuk earn star",
            "ChallengeStorage — JSON-backed SharedPrefs: scores, stars, streak, claimed rewards",
            "StreakTracker — anchor today/yesterday yang ada star, walk-back maksimum 365 hari",
            "3 tier reward: POWER_UP_3DAY, THEME_7DAY, XP_BONUS_14DAY",
            "CalendarView — grid 5×6 = 30 hari terakhir, status: empty / played / starred / today",
        ],
        page_num=11, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 12 — Fitur 6: Statistics & Achievements
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Fitur 6 — Statistics & Achievements",
        subtitle="20 achievements, per-mode statistics",
        bullets=[
            "20 achievements terbagi 5 kategori: SCORE (5), COMBO (3), LINES (3), GAMES (3), SPECIAL (6)",
            "Score: FIRST_GAME, CENTURION, THOUSAND_CLUB, FIVE_THOUSAND, TEN_THOUSAND",
            "Combo: COMBO_STARTER (×2), COMBO_MASTER (×5), COMBO_LEGEND (×10)",
            "Special: SPEED_DEMON, PERFECT, STREAK_3, STREAK_7, PRESTIGE, COMPLETIONIST (semua 19 lain)",
            "Per-mode statistics: games played, average score, best combo, total lines, play time, win streak",
            "Win streak threshold: skor > 1000 → streak +1; lainnya reset",
            "AchievementManager.evaluate*() — idempotent, persisted via SharedPreferences",
        ],
        page_num=12, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 13 — Fitur 7: Online Leaderboard (Firebase)
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Fitur 7 — Online Leaderboard (Firebase)",
        subtitle="Anonymous + Google Sign-In, snapshot listener live",
        bullets=[
            "Firebase Auth: anonymous default → bisa di-link ke Google account (preserve UID)",
            "Display name auto: 'PlayerNNNN' (4 digit) untuk anonymous user",
            "Firestore path: leaderboards/{mode}/scores/{userId}; doc ID = uid (1 score/mode/user)",
            "Top-100 per mode, all-time + weekly (Mon 00:00 UTC – Sun 23:59:59.999 UTC)",
            "Real-time updates via snapshot listener — UI auto-refresh tanpa pull",
            "Offline ScoreQueue: enqueue ke SharedPrefs JSON, sync on NetworkCallback.onAvailable",
            "Rate limit 10 detik client-side; Firestore rules validasi field set + score ∈ [0, 999999]",
        ],
        page_num=13, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 14 — Fitur 8: Level Progression & XP
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Fitur 8 — Level Progression & XP",
        subtitle="100 level, 7 reward, prestige loop",
        bullets=[
            "XP formula: floor((floor(score/10) + Σ(5·comboLevel) + 50·perfectClears) · (1 + 0.1·prestige))",
            "XP ke level berikutnya = currentLevel × 100",
            "100 level cap; multi-level-up didukung (loop sambil ada surplus XP)",
            "7 reward unlock: EXTRA_HINT (lv5), NEON_PALETTE (10), WOOD_SKIN (20), EXTRA_POWERUP_SLOT (30), RETRO_PALETTE (50), SPACE_SKIN (75), MASTER_BADGE (100)",
            "Prestige di lv 100 — reset ke lv 1, currentXP=0, prestigeCount+1, +10% XP multiplier permanen",
            "Reward yang sudah unlock dipertahankan setelah prestige",
            "Level-up overlay pop-up via GameCallback.onLevelUp(newLevels, newRewards)",
        ],
        page_num=14, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 15 — Fitur 9: Timed Mode
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Fitur 9 — Timed Mode",
        subtitle="Countdown HUD, multiplier band, time bonus",
        bullets=[
            "TimedModeController — pure JVM (depends only on android.os.Bundle)",
            "Countdown HUD warna: putih >30s, kuning >10s, merah ≤10s; pulse 2Hz di 5 detik terakhir",
            "Multiplier band: <15s=1.0×, <30s=1.5×, <45s=2.0×, <60s=2.5×, ≥60s=3.0× (timed90) / 2.5× (timed60)",
            "Time bonus +1.5s tiap line break, indicator flash 800ms",
            "End bonus: 5 poin per detik tersisa (END_BONUS_PER_SECOND)",
            "Auto-pause saat onPause() (Activity di-background)",
            "Restore selalu paused (Requirement 10.2) — user harus tap untuk resume",
        ],
        page_num=15, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 16 — Fitur 10: Themes & Customization
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Fitur 10 — Themes & Customization",
        subtitle="6 palette + 5 board skin, unlock by total score",
        bullets=[
            "6 ColorPalette: DEFAULT (gratis), NEON (5K), PASTEL (10K), RETRO (25K), DARK (50K), OCEAN (100K)",
            "5 BoardSkin: DEFAULT, WOOD (5K), METAL (25K), SPACE (50K), PIXEL_ART (100K)",
            "Skin paint via Canvas primitives saja — tanpa bitmap asset (gradient + hatching + procedural stars)",
            "Unlock berdasarkan totalScore akumulatif (semua mode dijumlahkan)",
            "ThemeManager singleton — load/persist palette+skin id, validate unlock, fallback ke DEFAULT",
            "SettingsActivity — custom Canvas view: palette grid 2×3, skin strip 1×5, preview 4×4",
            "Locked tile dim 50% black + lock glyph + threshold dalam emas",
        ],
        page_num=16, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 17 — Visual Rendering
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Visual Rendering",
        subtitle="100% Canvas, density-aware, pool reuse",
        bullets=[
            "Block beveled: 4 border colors (top/left/right/bottom) dari multiplier RGB Piece",
            "Density-aware layout — semua dimensi pakai dp × density px",
            "Ghost pulse alpha: computePulseAlpha(periodMs, min, max) berbasis SystemClock.uptimeMillis()",
            "ParticlePool fixed 512, drawCircle setiap frame; activeCount tracked",
            "GameView.calculateLayout() menyusun board, HUD, hand pieces, power-up slots responsively",
            "Pause overlay, game-over overlay, level-up overlay dirender on-top via FrameLayout",
        ],
        page_num=17, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 18 — Audio System
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Audio System",
        subtitle="SoundPool, 9 sfx, pitch modulation",
        bullets=[
            "AudioManager singleton — SoundPool dengan max 6 stream",
            "9 sound effect: place, break, gameover, juiciness, combo1..combo5",
            "Combo level→sound: idx = min(comboLevel-2, 4) → combo1..combo5",
            "playPerfectClear() — reuse juiciness dengan pitch 1.2×",
            "Power-up acquired → place.mp3 pitch 1.3×",
            "Power-up activate: BOMB→juiciness, LINE_SWEEP→break, ROTATE→place pitch 1.5×, UNDO→place pitch 0.7×",
            "Volume slider tersimpan di StorageManager (KEY_VOLUME)",
        ],
        page_num=18, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 19 — Persistence
    # ------------------------------------------------------------------
    make_two_col_slide(
        prs,
        title="Persistence",
        subtitle="Hybrid: SharedPreferences + Bundle + Firestore",
        left_title="Local (SharedPreferences)",
        left_bullets=[
            "tile_blast_prefs — high scores, volume, level/XP, prestige, unlocked rewards, total score, palette/skin id, statistics, achievements",
            "tile_blast_auth — userId, displayName, isAnonymous",
            "tile_blast_leaderboard — last submission, cached rank, last viewed mode",
            "tile_blast_score_queue — offline queued scores (JSON array)",
            "daily_challenge_prefs — daily scores, stars, streak, claimed rewards",
        ],
        right_title="Process (Bundle) + Cloud (Firestore)",
        right_bullets=[
            "Bundle: GameView.saveState/restoreState — board cells, hand, score, power-ups, snapshot, timer state",
            "Restore selalu paused (Timer & PowerUp activation reset ke IDLE)",
            "Firestore: leaderboards/{mode}/scores/{userId} dan players/{userId}",
            "Snapshot listener untuk live leaderboard",
            "ScoreQueue persist offline → sync on NetworkCallback",
            "BoardSnapshot.writeTo/readFrom Bundle — snapshot Undo survive rotation",
        ],
        page_num=19, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 20 — Property-Based Testing
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Property-Based Testing",
        subtitle="jqwik framework, ~30 properties",
        bullets=[
            "BoardSnapshot round-trip: capture → mutate → restore → equals(initial)",
            "Inventory cap invariant: count[type] selalu ≤ MAX_PER_TYPE = 2 setelah grant",
            "Milestone uniqueness: setiap milestone hanya award sekali per game",
            "Activation gating matrix: gating rule per (type × ctx) verified exhaustive",
            "BOMB apply preservation: only 3×3 area cleared, sisa board tidak berubah",
            "LINE_SWEEP apply preservation: hanya satu row atau column yang clear",
            "Rotate 4-cycle: applyRotate × 4 ≡ original matrix (modulo dimensi)",
            "Save/restore round trip: saveState(out) → restoreState(out) ≡ initial state",
        ],
        page_num=20, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 21 — Spec-Driven Development
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Spec-Driven Development",
        subtitle="11 specs di .kiro/specs/, requirements + design + tasks",
        bullets=[
            "Folder .kiro/specs/ berisi 11 spec: enhanced-visual-feedback, ghost-preview-hints, power-ups-system, special-tiles, daily-challenge, statistics-achievements, online-leaderboard-firebase, level-progression-xp, timed-mode, themes-customization, all-bugs-fix",
            "Tiap spec punya 3 dokumen: requirements.md, design.md, tasks.md",
            "Tasks dijalankan dalam waves yang bisa parallel (multi-subagent execution)",
            "Acceptance criteria → property test prework → property generation",
            "Bugfix workflow terpisah: bug condition methodology di all-bugs-fix",
        ],
        page_num=21, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 22 — Statistik Project
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Statistik Project",
        subtitle="Hitungan kasar dari source tree",
        bullets=[
            "50+ Java class di package com.allan.tileblast",
            "Source code total ±7000+ baris (game/ + effects/ + leaderboard/ + daily/ + …)",
            "10 fitur major terimplementasi end-to-end (UI + logic + persistence + tests)",
            "11 spec dengan tasks granular, BUILD SUCCESSFUL",
            "5 game mode dengan reuse 1 GameView class (parameterized via setup())",
            "20 achievement, 100 level cap, 6 palette × 5 skin × 4 power-up types",
        ],
        page_num=22, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 23 — Tantangan Teknis
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Tantangan Teknis",
        subtitle="Bug yang ditemui & cara solve",
        bullets=[
            "Race condition di PieceSnapAnim: piece flag terhapus sebelum commit → fix dengan hasPendingPiece() / clearPiece() pair",
            "Firestore rules deployment: default deny-all → app PERMISSION_DENIED sampai rules ter-publish",
            "Parallel subagent edit conflict di GameView.java (file paling besar) → perlu serialisasi wave",
            "Density-aware layout tuning: rect size yang fixed di emulator pecah di device fisik",
            "Bundle restore: timer/power-up activation harus selalu paused/IDLE untuk avoid stale state",
            "Snapshot deep-copy untuk Undo: matrix piece di-copy per slot agar rotation preserved",
        ],
        page_num=23, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 24 — Demo Highlights
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Demo Highlights",
        subtitle="Screenshot placeholder — akan diganti saat presentasi",
        bullets=[
            "[Insert screenshot: main menu — logo, level/XP bar, 4 mode + leaderboard/settings]",
            "[Insert screenshot: gameplay classic — board 8×8 dengan ghost preview emas]",
            "[Insert screenshot: chaos dengan special tiles — frozen, locked, bomb dengan countdown]",
            "[Insert screenshot: daily challenge — kalender 30 hari + streak fire icon]",
            "[Insert screenshot: leaderboard — tab All-Time/Weekly, highlight gold pemain sendiri]",
            "[Insert screenshot: settings — palette grid 2×3, skin strip 1×5, preview area]",
        ],
        page_num=24, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 25 — Future Work
    # ------------------------------------------------------------------
    make_content_slide(
        prs,
        title="Future Work",
        subtitle="Yang mungkin dikembangkan berikutnya",
        bullets=[
            "Multiplayer mode — head-to-head atau co-op via Firebase Realtime Database",
            "Custom themes editor — user bikin palette sendiri",
            "Replay system — record dan playback game dari seed + sequence",
            "Cloud save sync — backup progression cross-device via Firestore",
            "More piece shapes — 4-block, 5-block, atau special pieces dengan effect",
            "Online tournaments — bracket weekly dengan reward exclusive",
            "Accessibility audit — screen reader support, contrast checker",
        ],
        page_num=25, total=TOTAL,
    )

    # ------------------------------------------------------------------
    # Slide 26 — Penutup
    # ------------------------------------------------------------------
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_text_box(
        s, Inches(0.5), Inches(2.4),
        Inches(12.3), Inches(1.4),
        "Terima Kasih",
        size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER,
    )
    add_text_box(
        s, Inches(0.5), Inches(3.6),
        Inches(12.3), Inches(0.6),
        "Q & A",
        size=32, color=WHITE, align=PP_ALIGN.CENTER,
    )
    add_text_box(
        s, Inches(0.5), Inches(4.5),
        Inches(12.3), Inches(0.5),
        "Repository: [Insert Link Repo]",
        size=18, color=GREY, align=PP_ALIGN.CENTER,
    )
    add_text_box(
        s, Inches(0.5), Inches(5.0),
        Inches(12.3), Inches(0.5),
        "Kontak: [Insert Email / Handle]",
        size=14, color=GREY, align=PP_ALIGN.CENTER,
    )
    add_footer(s, 26, TOTAL)

    # ------------------------------------------------------------------
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
